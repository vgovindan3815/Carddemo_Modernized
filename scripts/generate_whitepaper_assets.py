from pathlib import Path
from datetime import date

from docx import Document
from docx.shared import Pt, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "deliverables"
OUT_DIR.mkdir(parents=True, exist_ok=True)

DOCX_PATH = OUT_DIR / "AI_Assisted_Mainframe_Reinvention_Whitepaper.docx"
PPTX_PATH = OUT_DIR / "AI_Assisted_Mainframe_Reinvention_Executive_Deck.pptx"

TODAY = date.today().isoformat()

ONLINE_TIMELINE = {
    "reverse_engineering": "0.25 days (~2 hours)",
    "architecture_specs": "0.25 days (~2 hours)",
    "core_rebuild": "0.50 days (~4 hours)",
    "enhancement_hardening": "0.50 days (~4 hours)",
    "total": "1.50 days (~12 hours, ~0.3 work weeks)",
}


def add_heading_style(doc: Document):
    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(11)


def add_title(doc: Document, title: str, subtitle: str):
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(title)
    run.bold = True
    run.font.size = Pt(22)

    p2 = doc.add_paragraph()
    p2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run2 = p2.add_run(subtitle)
    run2.italic = True
    run2.font.size = Pt(12)

    doc.add_paragraph()


def add_section(doc: Document, title: str, bullets=None, paragraphs=None):
    doc.add_heading(title, level=1)
    if paragraphs:
        for para in paragraphs:
            doc.add_paragraph(para)
    if bullets:
        for bullet in bullets:
            doc.add_paragraph(bullet, style="List Bullet")


def build_docx():
    doc = Document()
    add_heading_style(doc)

    add_title(
        doc,
        "Reinventing Mainframe Modernization with AI-Assisted Reverse + Forward Engineering",
        f"Publish-ready white paper | Generated on {TODAY}",
    )

    add_section(
        doc,
        "Attribution",
        paragraphs=[
            "This case study uses the AWS CardDemo legacy reference assets from the AWS Mainframe Modernization sample application set. "
            "We acknowledge AWS for providing a realistic enterprise baseline that enabled transparent reverse engineering, architecture mapping, and modernization validation.",
        ],
    )

    add_section(
        doc,
        "Executive Summary",
        paragraphs=[
            "Traditional modernization approaches—tool-led translation and rehosting—reduce immediate migration risk but often preserve legacy constraints in new runtime environments.",
            "This paper presents an alternative: AI-assisted reverse engineering plus forward engineering. Legacy intent is extracted from COBOL/JCL/copybooks and transformed into modern, API-first architecture with maintainable code, hardened security controls, and faster enhancement delivery.",
            "In this program, parity and advanced enhancements were delivered within one compressed modernization stream, where traditional delivery models typically require multiple months and phase gates.",
        ],
    )

    add_section(
        doc,
        "Why Legacy-First Modernization Is No Longer Enough",
        bullets=[
            "Translation-heavy approaches can keep legacy control flow and coupling intact.",
            "Generated code can be hard to maintain for mainstream engineering teams.",
            "Rehosting often shifts infrastructure, not architecture or operating model.",
            "Product innovation is delayed into later phases rather than delivered during modernization.",
        ],
    )

    add_section(
        doc,
        "Approach Used in This Engagement",
        bullets=[
            "Reverse engineering from source artifacts: COBOL, copybooks, JCL/proc/control files, scheduler metadata.",
            "Behavior extraction into explicit contracts: routes, validations, error semantics, data model, execution rules.",
            "Forward engineering into modern stack: Angular frontend, Node/Express APIs, SQLite operational model.",
            "Rapid enhancement loops (‘vibe coding’) with continuous build verification and UX hardening.",
        ],
    )

    add_section(
        doc,
        "Architecture Evolution",
        paragraphs=[
            "Legacy Pattern (simplified): Terminal UI -> CICS programs -> VSAM/DB2 with batch orchestration via JCL and scheduler chains.",
            "Modernized Pattern: Angular UI -> REST API layer (Express) -> relational persistence + role-based security + structured observability + OpenAPI/Swagger contract surface.",
            "Batch modernization remains API-driven and operationally explicit (job metadata, run history, artifacts, logs).",
        ],
    )

    add_section(
        doc,
        "Enhancements Delivered Beyond Functional Parity",
        bullets=[
            "Admin list workflows with Search + All actions and paging for cards and transactions.",
            "Detail-first account actions with reliable first-click scroll to account detail section.",
            "In-page Back navigation on detail views (no browser-back dependency).",
            "High Contrast accessibility mode with persisted user preference.",
            "Deterministic API error model and field-level validation handling.",
            "Swagger/OpenAPI docs exposed at /api-docs and /api-docs.json.",
            "Consistent modernized UX design system and responsive page structure.",
        ],
    )

    add_section(
        doc,
        "Security and Engineering Quality Outcomes",
        bullets=[
            "Session-based authentication with role checks for admin workflows.",
            "Server-side validation with consistent error contract and correlation IDs.",
            "Secure middleware stack and explicit API boundaries.",
            "Readable, maintainable code aligned to modern engineering practices.",
            "Developer onboarding acceleration via API documentation and runbook updates.",
        ],
    )

    add_section(
        doc,
        "Timeline and Speed-to-Value",
        paragraphs=[
            "Traditional modernization of this complexity (parity plus UX/security/operational enhancements) is typically executed in multiple stages across several months.",
            "In this engagement, reverse engineering, forward engineering, and advanced enhancements were completed in one compressed stream with rapid iterative validation.",
            "The timeline below is transparent for ONLINE scope only and reflects observed active implementation effort in this engagement.",
        ],
        bullets=[
            f"Legacy analysis + reverse engineering: {ONLINE_TIMELINE['reverse_engineering']}",
            f"Target architecture + specification baseline: {ONLINE_TIMELINE['architecture_specs']}",
            f"Core implementation (UI/API/DB): {ONLINE_TIMELINE['core_rebuild']}",
            f"Enhancement wave + hardening: {ONLINE_TIMELINE['enhancement_hardening']}",
            f"Total end-to-end transformation: {ONLINE_TIMELINE['total']}",
        ],
    )

    add_section(
        doc,
        "Business Impact",
        bullets=[
            "Reduced time-to-value by delivering innovation during migration, not after.",
            "Lower long-term maintenance risk through modern modular architecture.",
            "Improved security posture and operational transparency from day one.",
            "Expanded talent pool by reducing dependence on niche runtime/tool skills.",
        ],
    )

    add_section(
        doc,
        "Conclusion",
        paragraphs=[
            "Modernization outcomes should be evaluated by agility, maintainability, security, and product velocity—not only by whether legacy code still executes on new infrastructure.",
            "This case demonstrates that AI-assisted reverse + forward engineering can deliver a superior modernization outcome faster than translation/rehosting-only paths for organizations targeting long-term reinvention.",
        ],
    )

    doc.save(DOCX_PATH)


def _set_title(slide, text):
    title = slide.shapes.title
    if title:
        title.text = text
        for p in title.text_frame.paragraphs:
            p.font.size = PPt(32)
            p.font.bold = True


def _add_bullets(slide, left, top, width, height, heading, bullets):
    box = slide.shapes.add_textbox(PInches(left), PInches(top), PInches(width), PInches(height))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = heading
    p.font.bold = True
    p.font.size = PPt(20)
    p.space_after = PPt(8)

    for item in bullets:
        bp = tf.add_paragraph()
        bp.text = item
        bp.level = 0
        bp.font.size = PPt(16)


def build_pptx():
    prs = Presentation()

    # Slide 1
    s1 = prs.slides.add_slide(prs.slide_layouts[0])
    s1.shapes.title.text = "AI-Assisted Mainframe Reinvention"
    s1.placeholders[1].text = (
        "Why translation/rehosting-only modernization is outdated\n"
        "CardDemo case study (AWS legacy baseline)"
    )

    # Slide 2
    s2 = prs.slides.add_slide(prs.slide_layouts[5])
    _set_title(s2, "Attribution and Context")
    _add_bullets(
        s2,
        0.6,
        1.5,
        12.2,
        4.8,
        "Acknowledgment",
        [
            "Legacy baseline from AWS CardDemo sample assets (AWS Mainframe Modernization).",
            "Objective: reverse engineer legacy intent and rebuild into modern architecture.",
            "Scope covered online and batch design with iterative enhancement delivery.",
        ],
    )

    # Slide 3
    s3 = prs.slides.add_slide(prs.slide_layouts[5])
    _set_title(s3, "Why Traditional Paths Underperform")
    _add_bullets(
        s3,
        0.6,
        1.5,
        12.2,
        5.2,
        "Observed limitations",
        [
            "Tool-led translation often preserves legacy coupling and control flow.",
            "Rehosting shifts runtime but frequently defers architecture modernization.",
            "Generated/transformed code can be harder to maintain and evolve.",
            "Business innovation is delayed into later phases (months).",
        ],
    )

    # Slide 4
    s4 = prs.slides.add_slide(prs.slide_layouts[5])
    _set_title(s4, "Method Used: Reverse + Forward Engineering")
    _add_bullets(
        s4,
        0.6,
        1.5,
        12.2,
        5.2,
        "Execution model",
        [
            "Extract behavior and rules from COBOL/JCL/copybooks/scheduler assets.",
            "Define explicit contracts: APIs, validations, data model, UX behavior.",
            "Implement modern stack: Angular + Node/Express + SQLite.",
            "Apply rapid enhancement loops with continuous build and quality checks.",
        ],
    )

    # Slide 5
    s5 = prs.slides.add_slide(prs.slide_layouts[5])
    _set_title(s5, "Architecture: Before vs After")
    _add_bullets(
        s5,
        0.6,
        1.5,
        5.9,
        5.2,
        "Before",
        [
            "Terminal/CICS-centric flows",
            "Implicit contracts across programs/jobs",
            "Opaque operational dependencies",
        ],
    )
    _add_bullets(
        s5,
        6.7,
        1.5,
        5.9,
        5.2,
        "After",
        [
            "Web-native UX + REST APIs",
            "Role-based security + deterministic validation",
            "Observable runtime + OpenAPI/Swagger docs",
        ],
    )

    # Slide 6
    s6 = prs.slides.add_slide(prs.slide_layouts[5])
    _set_title(s6, "Enhancements Delivered in Primary Wave")
    _add_bullets(
        s6,
        0.6,
        1.5,
        12.2,
        5.2,
        "Beyond parity",
        [
            "Admin Search + All + paging for cards/transactions",
            "Reliable first-click scroll to account detail",
            "In-page Back actions on detail screens",
            "High Contrast accessibility mode with persisted preference",
            "Swagger at /api-docs and /api-docs.json",
        ],
    )

    # Slide 7
    s7 = prs.slides.add_slide(prs.slide_layouts[5])
    _set_title(s7, "Timeline Compression")
    _add_bullets(
        s7,
        0.6,
        1.5,
        12.2,
        5.4,
        "Observed ONLINE implementation effort",
        [
            f"Legacy reverse engineering: {ONLINE_TIMELINE['reverse_engineering']}",
            f"Architecture/spec finalization: {ONLINE_TIMELINE['architecture_specs']}",
            f"Core forward engineering: {ONLINE_TIMELINE['core_rebuild']}",
            f"Enhancements + hardening: {ONLINE_TIMELINE['enhancement_hardening']}",
            f"Total end-to-end: {ONLINE_TIMELINE['total']}",
            "Traditional equivalent: typically multiple months across phases",
            "Transparency: figures represent ONLINE stream only",
        ],
    )

    # Slide 8
    s8 = prs.slides.add_slide(prs.slide_layouts[5])
    _set_title(s8, "Security and Maintainability Outcomes")
    _add_bullets(
        s8,
        0.6,
        1.5,
        12.2,
        5.2,
        "Engineering quality",
        [
            "Session auth + role checks + explicit API boundaries",
            "Deterministic error model and field-level validation mapping",
            "Correlation IDs and structured logging",
            "Readable modular code with mainstream maintainability",
            "Faster onboarding using OpenAPI and updated runbooks",
        ],
    )

    # Slide 9
    s9 = prs.slides.add_slide(prs.slide_layouts[5])
    _set_title(s9, "Business Value")
    _add_bullets(
        s9,
        0.6,
        1.5,
        12.2,
        5.2,
        "Executive impact",
        [
            "Innovation delivered during modernization, not postponed",
            "Reduced long-term maintenance and skills risk",
            "Accelerated time-to-value and decision confidence",
            "Modern platform ready for further domain expansion",
        ],
    )

    # Slide 10
    s10 = prs.slides.add_slide(prs.slide_layouts[5])
    _set_title(s10, "Recommendation")
    _add_bullets(
        s10,
        0.6,
        1.7,
        12.2,
        4.8,
        "Decision guidance",
        [
            "Use translation/rehosting as tactical bridges, not end state.",
            "Adopt AI-assisted reverse + forward engineering for strategic reinvention.",
            "Measure success by agility, maintainability, security, and speed to enhancement.",
        ],
    )

    prs.save(PPTX_PATH)


def main():
    build_docx()
    build_pptx()
    print(f"Created: {DOCX_PATH}")
    print(f"Created: {PPTX_PATH}")


if __name__ == "__main__":
    main()
